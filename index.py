from flask import Flask, Response, render_template, flash, request, send_file,send_from_directory, url_for,redirect
from selenium.webdriver import Firefox
from selenium.webdriver.firefox.options import Options
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from urllib.request import urlopen, urlretrieve
from random import randint
from io import StringIO, BytesIO
from os import path
from anuncio import Anuncio

app = Flask(__name__)
app.secret_key = b'ADSREPORT_ASD3AD0AS_5#y2L"F4Q8z\n\xec]/'
anuncios = []
busqueda = ""
mes = ""
# Routes to Render Something
@app.route('/')
def home():
    return render_template("home.html")


@app.route('/download')
def download():
    flash("Información obtenida exitosamente...",'info')
    return render_template("download.html")


@app.route("/", methods = ['POST', 'GET'])
def enviarM():
    if request.method == 'POST':
        iddato = request.form['idbuscar']
        anuncios.clear()
        ObtenerData(iddato)
        return  redirect(url_for('download'))
        
def ObtenerData(dato):
    global busqueda, mes
    busqueda = dato
    opts = Options()
    opts.set_headless()
    assert opts.headless  # Operating in headless mode
    diract = path.dirname(path.abspath(__file__)) #Directorio donde se encuentra el script
    browser = Firefox(executable_path=diract+"/geckodriver.exe", options=opts)
    urlgetx = "https://www.facebook.com/ads/library/?active_status=all&ad_type=all&country=GT&q="+dato+"&sort_data[direction]=desc&sort_data[mode]=relevancy_monthly_grouped&search_type=keyword_unordered"
    browser.get(urlgetx)
    time.sleep(3)
    #------Listas
    element = browser.find_elements_by_class_name('_9ccv')
    for e in element:
        if e != element[0]:
            browser.execute_script("""
        var element = arguments[0];
        element.parentNode.removeChild(element);
        """, e)

    card = browser.find_elements_by_class_name('_99s5') 
    cn = 0
    for c in card:
        tipo = "Imagen"
        try: 
            fs = c.find_element_by_class_name('_9cd3').text
        except: pass
        try: 
            aus = c.find_element_by_class_name('_8nqq')
            ausi = aus.get_attribute('src')
            aust = aus.get_attribute('alt')
        except: 
            ausi = ""
            aust = "No data"
            pass
        try: 
            ds = c.find_elements_by_class_name('_4ik4')
            ds = ds[1].text
        except: 
            ds = "No hay Descripción"
            pass
        try: 
            cs = c.find_element_by_class_name('_7jys').get_attribute('src')
            tipo = "Imagen"
        except: pass
        try: 
            cs = c.find_element_by_tag_name('video').get_attribute('src')
            tipo = "Video"
        except: pass
        obja = Anuncio(tipo,cs,ds,fs,ausi,aust) 
        anuncios.append(obja)

    mes = browser.find_element_by_css_selector("._99s9 .l61y9joe").text
    browser.close()
    return 'a'


def DescargarArchivo(tipo, url, width):
    img_data  = ""
    if tipo == "imagen":
        try:
            img_data = BytesIO(urlopen(url).read())
        except:
            img_data = ""
            print("Un error ocurrio con la URL seleccionada")
    elif tipo == "video":


        try:
            img_data = BytesIO(urlopen(url).read())
        except:
            img_data = ""
            print("Un error ocurrio con la URL seleccionada")
    return img_data 
    
@app.route('/pptx', methods = ['POST', 'GET'])
def CrearPresentacion():
    if len(anuncios) > 1:
        prs=Presentation()
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)
        title=slide.shapes.title # assigning a title
        subtitle=slide.placeholders[1] # placeholder for subtitle
        title.text="Reporte de Facebook Ads Library para '"+busqueda+"'"
        subtitle.text=mes+"- Generado automaticamente por AdReport v1.0"
        for i in anuncios:    
            blank_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_slide_layout)
            title=slide.shapes.title
            subtitle=slide.placeholders[1] 
            title.text=" "
            subtitle.text=" "
            fecha = slide.shapes.add_textbox(Inches(1),Inches(0.2),width=Inches(2), height=Inches(0.5))
            fecha.text = i.getFecha()
            desc = slide.shapes.add_textbox(Inches(1),Inches(0.7),width=Inches(8.8), height=Inches(1.19))
            p = desc.text_frame
            para = p.add_paragraph()
            p.word_wrap = True
            para.text = i.getDescription()
            para.font.size = Pt(12)
            
            if i.getTipo() == "Imagen":
                path = DescargarArchivo("imagen",i.getMedia(),-1)
                left=Inches(1)
                top=Inches(2)
                try:
                    img=slide.shapes.add_picture(path,left,top,width=Inches(5), height=Inches(5))
                except:
                    print("ERROR")
                    pass
            elif i.getTipo() == "Video":
                path = DescargarArchivo("video",i.getMedia(),-1)
                left=Inches(1)
                top=Inches(2)
                try:
                    img=slide.shapes.add_movie(path,left,top,width=Inches(5), height=Inches(5),poster_frame_image=None, mime_type='video/mp4')
                except: pass

            pathAuth = DescargarArchivo("imagen",i.getAuthMedia(),-1)
            left=Inches(6.6)
            top=Inches(4)
            try:
                desca = slide.shapes.add_textbox(Inches(6.5),Inches(3.5),width=Inches(1), height=Inches(1)).text = "Anunciante:"
                desc = slide.shapes.add_textbox(Inches(7.9),Inches(4),width=Inches(1), height=Inches(1))
                img=slide.shapes.add_picture(pathAuth,left,top,width=Inches(1), height=Inches(1))
                p = desc.text_frame
                para = p.add_paragraph()
                p.word_wrap = True
                para.text = i.getAuthTexto()
                para.font.size = Pt(12)
            except:
                print("ERROR")
                pass
        file_name = 'AdsReport.pptx'
        out_file = StringIO()
        mem = BytesIO()
        mem.write(out_file.getvalue().encode())
        prs.save(mem)
        # seeking was necessary. Python 3.5.2, Flask 0.12.2
        mem.seek(0)
        out_file.close()    
        anuncios.clear()
        return send_file(
            mem,
            as_attachment=True,
            attachment_filename=file_name,
            mimetype='pptx'
        )

    
# Make sure this we are executing this file
if __name__ == '__main__':
    app.run(debug=True)
