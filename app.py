from selenium.webdriver import Firefox
from selenium.webdriver.firefox.options import Options
from openpyxl import load_workbook, Workbook
import os
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from urllib.request import urlopen, urlretrieve
from random import randint
from io import StringIO, BytesIO
from PIL import ImageFile
from pytube import YouTube
from assets.anuncio import Anuncio
diract = os.path.dirname(os.path.abspath(__file__)) #Directorio donde se encuentra el script

def MensajeI():
    textini = """
      ___       _ ______                           _   
     / _ \     | || ___ \                         | |  
    / /_\ \  __| || |_/ / ___  _ __    ___   _ __ | |_ 
    |  _  | / _` ||    / / _ \| '_ \  / _ \ | '__|| __|
    | | | || (_| || |\ \|  __/| |_) || (_) || |   | |_ 
    \_| |_/ \__,_|\_| \_|\___|| .__/  \___/ |_|    \__| v 1.2
                            | |                      
                            |_|                      
    """
    os.system('cls')
    print ("\n*********************************************************************")
    print (textini) 
    print ("*********************************************************************\n")


def MenuP():
    MensajeI()
    print("Ingresa la opcion solicitada: ")
    print("1) Descargar data desde Facebook Ad Library (Archivo Excel) ")
    print("2) Generar presentacion desde un Archivo de Excel")
    txt = input()
    if txt == "1": 
        print("Ingresa el termino a buscar")
        exc = input()
        ObtenerDataFacebook(exc)
    elif txt == "2": 
        print("Ingresa el nombre del archivo excel (este debe estar dentro de la carpeta, recuerda la extension .xlsx)")
        exc = input()
        DescargarExcel(exc)

def getsizes(uri):
    file = urlopen(uri)
    size = file.headers.get("content-length")
    if size: size = int(size)
    p = ImageFile.Parser()
    while 1:
        data = file.read(1024)
        if not data:
            break
        p.feed(data)
        if p.image:
            return size, p.image.size
            break
    file.close()
    return size, None

def DescargarArchivo(tipo, url, width):
    img_data  = ""
    if tipo == "yt":
        yt = YouTube(url)
        out_file = yt.streams.first().download()
        img_data = out_file
        try:
            pass
        except:
            img_data = ""
            print("Un error ocurrio con la URL seleccionadayt")
    elif tipo == "imagen":
        try:
            img_data = BytesIO(urlopen(url).read())
        except:
            img_data = ""
            print("Un error ocurrio con la URL seleccionada")
    elif tipo == "video":
        try:
            img_data = BytesIO(urlopen(url).read())
        except:
            img_data = ""
            print("Un error ocurrio con la URL seleccionada")
    return img_data 

def DescargarExcel(archivo):
    try:
        workbook = load_workbook(filename=diract+"/"+archivo)
        hoja = workbook.active
        all_rows = hoja.rows
        prs=Presentation()
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)
        title=slide.shapes.title 
        subtitle=slide.placeholders[1] 
        subtitle.text="- Generado automaticamente por AdReport v1.0"
        for row in hoja.iter_rows(values_only=True):
            own = row[0]
            img = row[1]
            blank_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_slide_layout)
            title=slide.shapes.title
            subtitle=slide.placeholders[1] 
            title.text=" "
            subtitle.text=" "
            anun = slide.shapes.add_textbox(Inches(1),Inches(0.2),width=Inches(4), height=Inches(0.5))
            p = anun.text_frame
            para = p.add_paragraph()
            p.word_wrap = True
            para.text = own
            para.font.size = Pt(40)
            left=Inches(1)
            top=Inches(1.2)
            if "https://youtube.com" in img:
                try:
                    print("Descargando video de Youtube:"+img)
                    path = DescargarArchivo("yt",img,-1)
                    imga=slide.shapes.add_movie(path,left,top,width=Inches(5), height=Inches(5),poster_frame_image=None, mime_type='video/mp4')
                    os.remove(path) ##Borrar video descargada, borrar esta linea de codigo si desea conservar la imagen descargada
                except: pass
            else:
                try:
                    print("Descargando Imagen de:"+img)
                    if not 'https:' in img:
                        img = "https:"+img
                    try:
                        path = DescargarArchivo("imagen",img,-1)
                        a = getsizes("https:"+img)[1][0] / 96
                        b = getsizes("https:"+img)[1][1] / 96
                        imga=slide.shapes.add_picture(path,left,top,width=Inches(a), height=Inches(b))
                    except: imga=slide.shapes.add_picture(path,left,top,width=Inches(5), height=Inches(5))
                except Exception as e:
                    try:
                        print("Descargando video de:"+img)
                        path = DescargarArchivo("video",img,-1)
                        img=slide.shapes.add_movie(path,left,top,width=Inches(5), height=Inches(5),poster_frame_image=None, mime_type='video/mp4')
                    except:
                            print("Ocurrio un error al descargar el siguiente video: "+img)
                except:
                    print("Ocurrio un error al descargar el siguiente archivo: "+img)
        print("Se guardó la presentación en el siguiente directorio: "+diract+"/AdReport.pptx")        
        prs.save(diract+"/AdReport.pptx") # saving file
    except: print("Ocurrio un error al intentar descargar el archivo excel, intentalo nuevamente")


anuncios = []
def ObtenerDataFacebook(dato):
    busqueda = dato
    opts = Options()
    opts.headless = True
    diract = os.path.dirname(os.path.abspath(__file__)) #Directorio donde se encuentra el script
    print("Obteniendo informacion para "+busqueda)        
    browser = Firefox(executable_path=diract+"/assets/geckodriver.exe", options=opts)
    urlgetx = "https://www.facebook.com/ads/library/?active_status=all&ad_type=all&country=GT&q="+dato+"&sort_data[direction]=desc&sort_data[mode]=relevancy_monthly_grouped&search_type=keyword_unordered"
    browser.get(urlgetx)
    time.sleep(3)
    #------Listas
    element = browser.find_elements_by_class_name('_9ccv')
    for e in element:
        if e != element[0]:
            browser.execute_script("""
        var element = arguments[0];
        element.parentNode.removeChild(element);
        """, e)
    card = browser.find_elements_by_class_name('_99s5') 
    cn = 0
    for c in card:
        tipo = "Imagen"
        try: 
            fs = c.find_element_by_class_name('_9cd3').text
        except: pass
        try: 
            aus = c.find_element_by_class_name('_8nqq')
            ausi = aus.get_attribute('src')
            aust = aus.get_attribute('alt')
        except: 
            ausi = ""
            aust = "No data"
            pass
        try: 
            ds = c.find_elements_by_class_name('_4ik4')
            ds = ds[1].text
        except: 
            ds = "No hay Descripción"
            pass
        try: 
            cs = c.find_element_by_class_name('_7jys').get_attribute('src')
            tipo = "Imagen"
        except: pass
        try: 
            cs = c.find_element_by_tag_name('video').get_attribute('src')
            tipo = "Video"
        except: pass
        obja = Anuncio(tipo,cs,ds,fs,ausi,aust) 
        anuncios.append(obja)
    mes = browser.find_element_by_css_selector("._99s9 .l61y9joe").text
    browser.close()
    print("Informacion obtenida exitosamente ")        
    print("Generando archivo Excel") 
    if os.path.exists(diract+"\AdsReport.xlsx"):
        wb = load_workbook(diract+"\AdsReport.xlsx")
        ws = wb.create_sheet(title=dato)
    else: 
        wb = Workbook()
        ws = wb.active
        ws.title = dato

    cn=1
    for i in anuncios:    
        fecha = i.getFecha()
        desc = i.getDescription()
        img = i.getMedia()
        auth = i.getAuthTexto()
        ws['A'+str(cn)] = auth
        ws['B'+str(cn)] = img
        ws['C'+str(cn)] = " " # Esta columna se agrea por motivos de orden únicamente
        ws['D'+str(cn)] = fecha
        cn+=1
    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 20
    ws.column_dimensions["C"].width = 5
    ws.column_dimensions["D"].width = 30
    wb.save(filename = diract+'/AdsReport.xlsx')
    print("Archivo de Excel generado exitosamente en: "+diract+'/AdsReport.xlsx') 

if __name__ == '__main__':
    MenuP()

